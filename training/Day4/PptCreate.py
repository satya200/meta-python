
from pptx import Presentation
import pptx.chart.data
from pptx.chart.data import ChartData
from pptx.util import Inches
from pptx.enum.chart import XL_CHART_TYPE

ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[5])
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'North']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
#print(dir(pptx.chart.data)

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
ppt.save("Mypresentation.pptx")


